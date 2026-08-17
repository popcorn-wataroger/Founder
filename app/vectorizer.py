"""ソースの本文抽出・ベクトル化を担う司令塔モジュール。

この段階では「① 本文抽出」のみ実装している。
embedding（ベクトル化）や Qdrant への保存は後続ステップで追加する。
"""

import io

import requests
from bs4 import BeautifulSoup
from docx import Document
from google import genai
from pptx import Presentation
from pypdf import PdfReader

# ファイルの読み出しは app/storage.py に任せる。
# 保存先がGCSかローカルかを、このモジュールは知らなくてよい
# （名前の検証も storage の中で必ず行われる）。
from app import storage
from app.config import GEMINI_API_KEY

# URLの安全性検証は app/safe_urls.py に集約している。
# このモジュールは「検証を通った安全なURL」を受け取って通信するだけにする
from app.safe_urls import URL_TIMEOUT, build_safe_public_url

# チャンク分割の設定（Issue #19 の完了条件）
CHUNK_SIZE = 500  # 1チャンクの最大文字数
CHUNK_OVERLAP = 100  # 隣り合うチャンク間で重ねる文字数

# embedding（ベクトル化）に使う Gemini のモデル名
# 後で「どのモデルを使ったか」を確認できるよう定数として持たせる
EMBEDDING_MODEL = "gemini-embedding-001"

# Gemini APIクライアント。APIキーは config 経由で読み込む（コードに直書きしない）
_client = genai.Client(api_key=GEMINI_API_KEY)


def extract_text(path: str, file_type: str) -> str:
    """ソースから本文テキストを取り出す。

    入力:
        path      … ストレージ上のファイルを指すキー（DBの file_path）。
                    ただし file_type が 'url' の場合はURL文字列。
        file_type … ソースの種別（'pdf' / 'docx' / 'pptx' / 'txt' / 'url'）

    出力:
        本文テキスト（文字列）

    file_type ごとに専用の抽出処理へ振り分ける（ディスパッチ）。
    未対応の形式が来たら ValueError を投げる。
    """
    if file_type == "pdf":
        return _extract_pdf(path)
    if file_type == "docx":
        return _extract_docx(path)
    if file_type == "pptx":
        return _extract_pptx(path)
    if file_type == "txt":
        return _extract_txt(path)
    if file_type == "url":
        return _extract_url(path)
    raise ValueError(f"未対応のファイル形式です: {file_type}")


def _read_file(path: str) -> io.BytesIO:
    """ストレージからファイルの中身を読み出し、ライブラリに渡せる形にして返す。

    入力:
        path … ストレージ上のファイルを指すキー（DBの file_path）

    出力:
        中身を載せた BytesIO（ファイルのように読めるオブジェクト）

    なぜパスではなく中身を渡すのか:
        保存先がGCSになると「ローカルのパス」は存在しない。
        PDF/Word/PowerPoint のライブラリはいずれもファイルのように読める
        オブジェクトを受け取れるので、中身をメモリ上に載せて渡す。
        これで、保存先がどこであっても抽出処理は同じコードで済む。

    安全性について:
        パスの検証は storage.read_bytes の中で必ず行われる
        （保存名の生成規則に合わない値は ValueError になる）。
        このモジュールが検証ルールを書き写す必要はない。
    """
    return io.BytesIO(storage.read_bytes(path))


def _extract_pdf(path: str) -> str:
    """PDFから本文を抽出する。"""
    reader = PdfReader(_read_file(path))
    # ページごとにテキストを取り出し、改行で連結する
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extract_docx(path: str) -> str:
    """Word(.docx)から本文を抽出する。"""
    document = Document(_read_file(path))
    # 段落（paragraph）ごとの文字列を改行で連結する
    paragraphs = [para.text for para in document.paragraphs]
    return "\n".join(paragraphs)


def _extract_pptx(path: str) -> str:
    """PowerPoint(.pptx)から本文を抽出する。"""
    presentation = Presentation(_read_file(path))
    texts: list[str] = []
    # スライド → 図形(shape) の順にたどり、テキストを持つ図形だけ集める
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _extract_txt(path: str) -> str:
    """テキスト(.txt)をそのまま読み込む。"""
    # 文字化けを避けるため UTF-8 として読む
    return storage.read_bytes(path).decode("utf-8")


def _extract_url(url: str) -> str:
    """URLのウェブページから本文を抽出する。

    入力:
        url … 取得したいウェブページのURL

    出力:
        ページの表示テキスト（文字列）

    処理（SSRF対策）:
        URLは社長がフォームから自由に入力できる。もし何の検査もせずに
        requests.get すると、サーバー自身に「内部だけに見えるアドレス」へ
        アクセスさせられてしまう（SSRF＝サーバー側リクエスト偽造）。
        たとえば http://169.254.169.254/ はクラウド(GCP等)のメタデータサーバーで、
        サービスアカウントのアクセストークンが取れてしまう。
        http://localhost:8000/ や社内ネットワークのIPも同様に危険。

        そこで build_safe_public_url（app/safe_urls.py）で検査済みの要素だけから
        URLを組み立て直し、入力された url そのものではなく、
        その safe_url だけを requests.get に渡す。
    """
    safe_url = build_safe_public_url(url)

    # sink に到達するのは、検証済み要素から作り直した safe_url だけ。
    # timeout … 応答が無いURLで固まらないための保険
    # allow_redirects=False … 公開URLに見せかけて内部アドレスへ302転送する
    #                         抜け道を塞ぐため、リダイレクトは追わない
    response = requests.get(safe_url, timeout=URL_TIMEOUT, allow_redirects=False)
    response.raise_for_status()

    # HTMLを解析し、本文に不要なタグ（script/style）を取り除く
    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()

    # 表示テキストだけを取り出す
    return soup.get_text(separator="\n", strip=True)


def split_into_chunks(text: str) -> list[str]:
    """本文テキストをチャンク（小さな塊）のリストに分割する。

    入力:
        text … 分割したい本文テキスト

    出力:
        チャンク文字列のリスト

    なぜ分割するか:
        長文をそのままベクトル化すると要点がぼやけ、検索精度が落ちる。
        500文字ごとの小さな塊に分けることで、質問に近い部分だけを探しやすくする。

    なぜ重ねる（オーバーラップ）か:
        区切り目で文が途中に切れると意味が失われる。
        前のチャンクの末尾100文字を次のチャンクの先頭に重ねることで、
        文脈が境界で途切れるのを防ぐ。
    """
    # 空文字なら分割対象がないので空リストを返す
    if not text:
        return []

    # 次のチャンク開始位置をどれだけ進めるか（500 - 100 = 400文字ずつ前進）
    step = CHUNK_SIZE - CHUNK_OVERLAP

    chunks: list[str] = []
    start = 0
    while start < len(text):
        # 現在位置から最大500文字を1チャンクとして切り出す
        end = start + CHUNK_SIZE
        chunks.append(text[start:end])
        # 末尾まで到達したら終了（重複した余分なチャンクを作らない）
        if end >= len(text):
            break
        # オーバーラップ分を残して開始位置を前進させる
        start += step

    return chunks


def embed_text(text: str) -> list[float]:
    """1チャンクのテキストを数値ベクトル（embedding）に変換する。

    入力:
        text … ベクトル化したいテキスト（通常はチャンク1つ）

    出力:
        数値（float）のリスト＝ベクトル。
        gemini-embedding-001 は既定で3072次元のベクトルを返す。
        実際の次元数は len(戻り値) で後から確認できる。

    処理:
        Gemini の embedding 用モデルにテキストを渡し、
        「意味」を表す数値の並びを受け取る。
        この数値の近さ同士を比べることで、後段のRAGで
        質問に意味が近いチャンクを検索できるようになる。
    """
    # Gemini にテキストを送ってベクトルを生成してもらう
    result = _client.models.embed_content(
        model=EMBEDDING_MODEL,
        contents=text,
    )
    # ベクトルが得られなかった場合、代わりの値（ゼロベクトル等）で埋めてはいけない。
    # 壊れたベクトルを保存すると、以後の検索が静かに誤った結果を返し続けるため、
    # ここで止めて呼び出し元（登録処理）のロールバックに任せる。
    # 例外メッセージには本文を含めない（社内文書の中身が流出しないようにするため）
    if not result.embeddings:
        raise RuntimeError(
            f"Gemini がベクトルを返しませんでした model={EMBEDDING_MODEL} 文字数={len(text)}"
        )

    # 戻り値は embedding のリスト。今回は1件なので先頭の数値列(values)を返す
    values = result.embeddings[0].values
    if values is None:
        raise RuntimeError(
            f"Gemini の応答にベクトルの数値列が含まれていません model={EMBEDDING_MODEL} "
            f"文字数={len(text)}"
        )

    return values
